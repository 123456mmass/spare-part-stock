"""
Generate presentation slides from the EGAT template.
Keep original backgrounds, shapes, and styling.
Only modify text and add feature images on slides 2-3.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TEMPLATE = os.path.join(BASE, "docs", "slice", "02 แบบฟอร์มการเสนอแนวคิดระดับฝ่าย_2569.pptx")
IMG_DIR = os.path.join(BASE, "docs", "images")
MOBILE_DIR = os.path.join(BASE, "docs", "mobile app image")
OUT_PATH = os.path.join(BASE, "docs", "slice", "SparePartStock_Presentation.pptx")

prs = Presentation(TEMPLATE)
slides = list(prs.slides)

# ══════════════════════════════════════════════════════════════
# SLIDE 1 — Cover page: fill in department name
# ══════════════════════════════════════════════════════════════
slide1 = slides[0]

for shape in slide1.shapes:
    if shape.has_text_frame:
        text = shape.text_frame.paragraphs[0].text
        # Find the "ระบุชื่อฝ่าย..." text box and replace with our department
        if "ระบุชื่อฝ่าย" in text:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = "ฝ่ายจัดซื้อและคลังอะไหล่"
            # Also clear other runs in that paragraph
            for para in shape.text_frame.paragraphs:
                if "ระบุชื่อฝ่าย" in para.text:
                    # Keep only first run with new text, clear rest
                    runs = para.runs
                    if len(runs) > 0:
                        runs[0].text = "ฝ่ายจัดซื้อและคลังอะไหล่"
                        for r in runs[1:]:
                            r.text = ""

# ══════════════════════════════════════════════════════════════
# SLIDE 2 — Idea details: fill in all fields
# ══════════════════════════════════════════════════════════════
slide2 = slides[1]

# Map of label text → new content to put in the freeform area below it
# We need to add text inside the freeform boxes (they act as content areas)
# Strategy: find each label, then add a new textbox in the content area below it

# The freeform shapes act as colored background boxes
# Text labels sit inside them
# We add new text boxes for content in the remaining space of each box

def find_shape_by_text(slide, search_text):
    """Find a shape containing specific text."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                if search_text in para.text:
                    return shape
    return None

def find_freeform_by_pos(slide, top_approx):
    """Find freeform shape near a given top position."""
    for shape in slide.shapes:
        if shape.shape_type == 5:  # FREEFORM
            if abs(shape.top - top_approx) < 200000:
                return shape
    return None

# --- Fill "ชื่อแนวคิด" section ---
name_label = find_shape_by_text(slide2, "ชื่อแนวคิด")
if name_label:
    # The freeform behind "ชื่อแนวคิด" is at top=2447582
    # Add text inside the freeform area
    freeform = find_freeform_by_pos(slide2, 2447582)
    if freeform:
        txBox = slide2.shapes.add_textbox(
            freeform.left + Emu(200000),
            name_label.top + Emu(500000),
            freeform.width - Emu(400000),
            freeform.height - Emu(600000)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "ระบบจัดการสต็อกอะไหล่ (Spare Part Stock Management System)"
        for run in p.runs:
            run.font.size = Pt(22)

# --- Fill "ที่มา/ปัญหา" section ---
origin_label = find_shape_by_text(slide2, "ที่มา/ปัญหา")
if origin_label:
    freeform = find_freeform_by_pos(slide2, 4218100)
    if freeform:
        txBox = slide2.shapes.add_textbox(
            freeform.left + Emu(200000),
            origin_label.top + Emu(500000),
            freeform.width - Emu(400000),
            freeform.height - Emu(700000)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = ("• ข้อมูลอะไหล่กระจัดกระจายในไฟล์ Excel หลายไฟล์ ไม่มีฐานข้อมูลกลาง\n"
                  "• การค้นหาอะไหล่ใช้เวลา 5–10 นาทีต่อรายการ\n"
                  "• ไม่สามารถตรวจสอบจำนวนคงเหลือแบบ Real-time ได้\n"
                  "• การเพิ่มอะไหล่ใหม่มีขั้นตอนมาก ต้องถ่ายรูป ย้ายไฟล์ แทรกรูปใน Excel\n"
                  "• เสี่ยงอะไหล่ซ้ำซ้อนหรือขาดอะไหล่สำคัญ")
        for run in p.runs:
            run.font.size = Pt(18)

# --- Fill "แนวทางในการดำเนินการ" section ---
approach_label = find_shape_by_text(slide2, "แนวทางในการดำเนินการ")
if approach_label:
    freeform = find_freeform_by_pos(slide2, 6994067)
    if freeform:
        txBox = slide2.shapes.add_textbox(
            freeform.left + Emu(200000),
            approach_label.top + Emu(500000),
            freeform.width - Emu(400000),
            freeform.height - Emu(700000)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = ("1. พัฒนาระบบเว็บแอปพลิเคชันด้วย Next.js + React + TypeScript\n"
                  "2. ใช้ Prisma ORM จัดการฐานข้อมูล SQLite\n"
                  "3. ใช้ AI Assistant ค้นหาอะไหล่ด้วยข้อความและรูปภาพ\n"
                  "4. พัฒนาแอปพลิเคชันมือถือ Android ด้วย Flutter\n"
                  "5. เชื่อมต่อ LINE ผ่าน LIFF และ LINE Bot\n"
                  "6. รองรับ QR Code / Barcode / นำเข้า Excel")
        for run in p.runs:
            run.font.size = Pt(18)

# --- Fill "คุณค่าที่ได้รับ" section ---
value_label = find_shape_by_text(slide2, "คุณค่าที่(คาดว่าจะ)ได้รับ")
if value_label:
    freeform = find_freeform_by_pos(slide2, 6955967)
    if freeform:
        # This freeform is the left-bottom one
        txBox = slide2.shapes.add_textbox(
            freeform.left + Emu(200000),
            value_label.top + Emu(500000),
            Emu(5000000),
            freeform.height - Emu(700000)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = ("• ลดเวลาค้นหาอะไหล่จาก 5–10 นาที → ไม่เกิน 1 นาที\n"
                  "• ตรวจสอบจำนวนคงเหลือแบบ Real-time\n"
                  "• ลดขั้นตอนการเพิ่มอะไหล่ใหม่ด้วย AI\n"
                  "• เข้าถึงข้อมูลได้ทุกช่องทาง")
        for run in p.runs:
            run.font.size = Pt(17)

# --- Fill "รายชื่อผู้เสนอแนวคิด" section ---
team_label = find_shape_by_text(slide2, "รายชื่อผู้เสนอแนวคิด")
if team_label:
    name_list_shape = find_shape_by_text(slide2, "1.    ชื่อ-นามสกุล")
    if name_list_shape:
        tf = name_list_shape.text_frame
        # Clear and rewrite
        for i, para in enumerate(tf.paragraphs):
            if i == 0:
                para.runs[0].text = "1.    นายธนวัฒน์ จันทร์แก้ว          650xxxxx          ฝ่ายจัดซื้อและคลังอะไหล่"
                for r in para.runs[1:]:
                    r.text = ""
            elif i == 1:
                para.runs[0].text = "2."
            elif i == 2:
                para.runs[0].text = "3."
            elif i == 3:
                para.runs[0].text = "4."
            elif i == 4:
                para.runs[0].text = "5."

# --- Fill "ชื่อทีม" section ---
team_name_label = find_shape_by_text(slide2, "ชื่อทีม")
if team_name_label:
    for para in team_name_label.text_frame.paragraphs:
        for run in para.runs:
            if "ชื่อทีม" in run.text:
                run.text = "ชื่อทีม : ทีมพัฒนาระบบจัดการสต็อกอะไหล่"
            elif "…" in run.text:
                run.text = ""

# --- Fill category header ---
cat_label = find_shape_by_text(slide2, "ประเภท Process Innovation")
if cat_label:
    for para in cat_label.text_frame.paragraphs:
        for run in para.runs:
            if "Process Innovation" in run.text:
                run.text = "ประเภท Process Innovation — ระบบจัดการสต็อกอะไหล่ด้วยเทคโนโลยีดิจิทัล"

# ══════════════════════════════════════════════════════════════
# SLIDE 3 — Feature showcase: 3 versions (Web, Mobile, LINE)
# ══════════════════════════════════════════════════════════════
slide3 = slides[2]

# Fill category header on slide 3
for shape in slide3.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            if "ประเภท Process Innovation" in para.text:
                for run in para.runs:
                    if "Process Innovation" in run.text:
                        run.text = "ฟีเจอร์เด่น — ระบบจัดการสต็อกอะไหล่ 3 ช่องทางเข้าถึง"

# The slide 3 has one large freeform (the content area)
# Layout: 3 columns for Web App | Mobile App | LINE App
# Slide width: 18288000 EMU (20 inches)
# Freeform: left=459263, width=17371537

content_left = Emu(600000)
content_top = Emu(1100000)
content_width = Emu(17000000)
content_height = Emu(8800000)

col_width = Emu(5400000)
col_gap = Emu(300000)
col_starts = [
    content_left,
    content_left + col_width + col_gap,
    content_left + 2 * (col_width + col_gap),
]

# --- Column 1: Web App ---
# Title
txBox = slide3.shapes.add_textbox(col_starts[0], content_top, col_width, Emu(500000))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "🌐 Web App"
run.font.size = Pt(28)
run.font.bold = True
run.font.color.rgb = RGBColor(0x1A, 0x73, 0xE8)

# Subtitle
txBox2 = slide3.shapes.add_textbox(col_starts[0], content_top + Emu(500000), col_width, Emu(350000))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run()
run2.text = "Next.js + React + TypeScript"
run2.font.size = Pt(16)
run2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

# Dashboard image
img_path = os.path.join(IMG_DIR, "01-dashboard.png")
if os.path.exists(img_path):
    slide3.shapes.add_picture(
        img_path,
        col_starts[0] + Emu(200000),
        content_top + Emu(900000),
        width=Emu(5000000)
    )

# Feature bullets
txBox3 = slide3.shapes.add_textbox(
    col_starts[0] + Emu(100000),
    content_top + Emu(4800000),
    col_width - Emu(200000),
    Emu(4000000)
)
tf3 = txBox3.text_frame
tf3.word_wrap = True
p3 = tf3.paragraphs[0]
run3 = p3.add_run()
run3.text = ("✅ Dashboard สรุปข้อมูลแบบ Real-time\n"
             "✅ ค้นหาอะไหล่ด้วย AI Assistant\n"
             "✅ สแกน QR Code / Barcode\n"
             "✅ นำเข้า–ส่งออก Excel\n"
             "✅ บันทึกการเคลื่อนไหวสต็อก\n"
             "✅ จัดการหมวดหมู่ อาคาร บล็อก")
run3.font.size = Pt(16)

# --- Column 2: Mobile App (Flutter) ---
txBox = slide3.shapes.add_textbox(col_starts[1], content_top, col_width, Emu(500000))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "📱 Mobile App"
run.font.size = Pt(28)
run.font.bold = True
run.font.color.rgb = RGBColor(0x0D, 0x96, 0x88)

txBox2 = slide3.shapes.add_textbox(col_starts[1], content_top + Emu(500000), col_width, Emu(350000))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run()
run2.text = "Flutter (Android APK)"
run2.font.size = Pt(16)
run2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

# Mobile app screenshot
mobile_img = os.path.join(MOBILE_DIR, "83919_0.jpg")
if os.path.exists(mobile_img):
    slide3.shapes.add_picture(
        mobile_img,
        col_starts[1] + Emu(1200000),
        content_top + Emu(900000),
        width=Emu(3000000)
    )

# Feature bullets
txBox3 = slide3.shapes.add_textbox(
    col_starts[1] + Emu(100000),
    content_top + Emu(4800000),
    col_width - Emu(200000),
    Emu(4000000)
)
tf3 = txBox3.text_frame
tf3.word_wrap = True
p3 = tf3.paragraphs[0]
run3 = p3.add_run()
run3.text = ("✅ ใช้งานบน Android โดยตรง\n"
             "✅ Bottom Nav 4 แท็บ\n"
             "    แดชบอร์ด / อะไหล่ / สแกน / เมนู\n"
             "✅ สแกนบาร์โค้ดด้วยกล้อง\n"
             "✅ เพิ่มอะไหล่ + ถ่ายรูป + AI แนะนำ\n"
             "✅ รับเข้า–เบิกออกสต็อกทันที")
run3.font.size = Pt(16)

# --- Column 3: LINE App ---
txBox = slide3.shapes.add_textbox(col_starts[2], content_top, col_width, Emu(500000))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "💬 LINE App"
run.font.size = Pt(28)
run.font.bold = True
run.font.color.rgb = RGBColor(0x06, 0xC7, 0x55)

txBox2 = slide3.shapes.add_textbox(col_starts[2], content_top + Emu(500000), col_width, Emu(350000))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run()
run2.text = "LINE LIFF + LINE Bot"
run2.font.size = Pt(16)
run2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

# LINE Bot image
line_img = os.path.join(IMG_DIR, "13-line-stock-summary.png")
if os.path.exists(line_img):
    slide3.shapes.add_picture(
        line_img,
        col_starts[2] + Emu(500000),
        content_top + Emu(900000),
        width=Emu(4400000)
    )

# Feature bullets
txBox3 = slide3.shapes.add_textbox(
    col_starts[2] + Emu(100000),
    content_top + Emu(4800000),
    col_width - Emu(200000),
    Emu(4000000)
)
tf3 = txBox3.text_frame
tf3.word_wrap = True
p3 = tf3.paragraphs[0]
run3 = p3.add_run()
run3.text = ("✅ LINE LIFF — เปิดระบบใน LINE\n"
             "    สแกน QR เพิ่มอะไหล่ได้ทันที\n"
             "✅ LINE Bot — แชทสอบถามสต็อก\n"
             "    สรุปตามหมวดหมู่ / ค้นหารายการ\n"
             "✅ รองรับแชทกลุ่ม\n"
             "✅ ตอบกลับด้วย Flex Message")
run3.font.size = Pt(16)

# ── Save ──
prs.save(OUT_PATH)
print(f"Slides saved to: {OUT_PATH}")
